import streamlit as st
import os
import tempfile
import zipfile
import io
import base64
from PIL import Image
import pandas as pd
import numpy as np
import math

# Check for optional dependencies and show warnings if missing
missing_deps = []

# For document processing
import PyPDF2
import pdfplumber
import docx
import markdown

# Try to import optional packages
try:
    from pptx import Presentation
    has_pptx = True
except ImportError:
    has_pptx = False
    missing_deps.append("python-pptx")

try:
    import bs4
    has_bs4 = True
except ImportError:
    has_bs4 = False
    missing_deps.append("beautifulsoup4")

try:
    import tabulate
    has_tabulate = True
except ImportError:
    has_tabulate = False
    missing_deps.append("tabulate")

# Try to import scikit-learn for improved chunking
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    has_sklearn = True
except ImportError:
    has_sklearn = False
    missing_deps.append("scikit-learn")

# For NLP
import nltk
import re
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.probability import FreqDist
from string import punctuation
from collections import Counter

# Initialize NLTK downloads
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords', quiet=True)


# Patch NLTK's PunktTokenizer to avoid punkt_tab issues
try:
    from nltk.tokenize.punkt import PunktTokenizer
    
    # Save the original load_lang method
    original_load_lang = PunktTokenizer.load_lang
    
    # Define a new load_lang method that doesn't fail
    def patched_load_lang(self, lang):
        try:
            # Try the original method first
            return original_load_lang(self, lang)
        except:
            # Fall back to a simple split if the original method fails
            from nltk.tokenize.punkt import PunktParameters
            self._params = PunktParameters()
            return
    
    # Apply the patch
    PunktTokenizer.load_lang = patched_load_lang
    print("Applied PunktTokenizer patch")
except:
    print("Could not patch PunktTokenizer")

# Set page configuration
st.set_page_config(
    page_title="FileMerger App",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Display warnings about missing dependencies
if missing_deps:
    st.warning(f"Some features may be limited. Missing dependencies: {', '.join(missing_deps)}. Run 'pip install {' '.join(missing_deps)}' to install them.")

# CSS for styling
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        margin-bottom: 1rem;
        color: #2c3e50;
    }
    .sub-header {
        font-size: 1.5rem;
        font-weight: semi-bold;
        margin-bottom: 1rem;
        color: #34495e;
    }
    .file-info {
        background-color: #f8f9fa;
        padding: 1rem;
        border-radius: 5px;
        margin-bottom: 1rem;
    }
    .sidebar-header {
        font-size: 1.2rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = []
if 'conversion_done' not in st.session_state:
    st.session_state.conversion_done = False
if 'processed_files' not in st.session_state:
    st.session_state.processed_files = []
if 'merged_content' not in st.session_state:
    st.session_state.merged_content = ""
if 'current_page' not in st.session_state:
    st.session_state.current_page = "home"

# Function to extract text from PDF
def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
            text += "\n\n"
    return text

# Function to extract text from DOCX
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from TXT
def extract_text_from_txt(file):
    return file.read().decode("utf-8")

# Function to extract text from CSV/XLSX
def extract_text_from_table(file, file_extension):
    if file_extension.lower() == '.csv':
        df = pd.read_csv(file)
    else:  # xlsx
        df = pd.read_excel(file)
    
    # Convert DataFrame to markdown table
    return df.to_markdown(index=False)

# Function to extract text from PowerPoint
def extract_text_from_pptx(file):
    if not has_pptx:
        return "Python-pptx package is required for PowerPoint processing. Please install it with: pip install python-pptx"
    
    try:
        prs = Presentation(file)
        text = ""
        
        # Get title and add as heading
        if prs.slide_layouts:
            text += "# PowerPoint Presentation\n\n"
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            text += f"## Slide {i+1}\n\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n\n"
            
            text += "\n"
        
        return text
    except Exception as e:
        return f"Error processing PowerPoint file: {str(e)}"

# Function to extract text from HTML
def extract_text_from_html(file):
    if not has_bs4:
        return "BeautifulSoup4 is required for HTML processing. Please install it with: pip install beautifulsoup4"
    
    try:
        from bs4 import BeautifulSoup
        html_content = file.read().decode("utf-8")
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Extract text and preserve some structure
        text = ""
        
        # Get title if it exists
        title = soup.find('title')
        if title:
            text += f"# {title.text}\n\n"
        
        # Get headings and paragraphs
        for elem in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li']):
            tag_name = elem.name
            
            if tag_name.startswith('h'):
                level = int(tag_name[1])
                prefix = '#' * level
                text += f"{prefix} {elem.text.strip()}\n\n"
            elif tag_name == 'p':
                text += f"{elem.text.strip()}\n\n"
            elif tag_name == 'li':
                text += f"- {elem.text.strip()}\n"
        
        return text
    except Exception as e:
        return f"Error processing HTML file: {str(e)}"

# Function to process files based on type
def process_file(file, file_extension):
    try:
        if file_extension.lower() in ['.pdf']:
            return extract_text_from_pdf(file)
        elif file_extension.lower() in ['.docx', '.doc']:
            return extract_text_from_docx(file)
        elif file_extension.lower() in ['.txt', '.md', '.text', '.markdown']:
            return extract_text_from_txt(file)
        elif file_extension.lower() in ['.csv', '.xlsx', '.xls']:
            return extract_text_from_table(file, file_extension.lower())
        elif file_extension.lower() in ['.pptx', '.ppt']:
            return extract_text_from_pptx(file)
        elif file_extension.lower() in ['.html', '.htm']:
            return extract_text_from_html(file)
        elif file_extension.lower() in ['.json', '.xml', '.yaml', '.yml']:
            return extract_text_from_txt(file)  # Simple text extraction for these formats
        else:
            return f"Unsupported file format: {file_extension}. File will be treated as plain text."
    except Exception as e:
        return f"Error processing file: {str(e)}"

# Function to chunk text based on paragraphs with smart merging
def chunk_by_paragraphs(text, max_paragraphs=5):
    """Improved paragraph chunking that tries to keep related paragraphs together."""
    paragraphs = [p for p in text.split('\n\n') if p.strip()]
    
    if not paragraphs:
        return [""]
    
    chunks = []
    current_chunk = []
    current_count = 0
    
    for i, para in enumerate(paragraphs):
        # Check if this paragraph starts with a heading
        is_heading = bool(re.match(r'^#{1,6}\s', para))
        
        # If we've reached max paragraphs and the next one is a heading, 
        # or we're at max_paragraphs, create a new chunk
        if (current_count >= max_paragraphs) or \
           (current_count > 0 and is_heading and current_count >= max(2, max_paragraphs // 2)):
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = [para]
            current_count = 1
        else:
            current_chunk.append(para)
            current_count += 1
    
    # Add any remaining paragraphs
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    return chunks

# Function to chunk text based on semantic similarity
def chunk_by_semantic_similarity(text, max_tokens=500, similarity_threshold=0.3):
    """Chunk text based on semantic similarity between paragraphs."""
    # Fall back to token-based chunking if scikit-learn is not available
    if not has_sklearn:
        return chunk_by_tokens(text, max_tokens)
    
    try:
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        
        if not paragraphs:
            return [""]
        
        # If we have very few paragraphs, just use token-based chunking
        if len(paragraphs) <= 3:
            return chunk_by_tokens(text, max_tokens)
        
        # Create TF-IDF vectors for paragraphs
        vectorizer = TfidfVectorizer(stop_words='english')
        try:
            tfidf_matrix = vectorizer.fit_transform(paragraphs)
            similarities = cosine_similarity(tfidf_matrix)
        except:
            # If vectorization fails, fall back to token chunking
            return chunk_by_tokens(text, max_tokens)
        
        # Create chunks based on semantic similarity
        chunks = []
        current_chunk = [paragraphs[0]]
        current_token_count = len(word_tokenize(paragraphs[0]))
        current_idx = 0
        
        for i in range(1, len(paragraphs)):
            paragraph_token_count = len(word_tokenize(paragraphs[i]))
            similarity_to_current = similarities[current_idx, i]
            
            # If this paragraph is similar to the current one and we're not over the token limit,
            # add it to the current chunk
            if similarity_to_current > similarity_threshold and current_token_count + paragraph_token_count <= max_tokens:
                current_chunk.append(paragraphs[i])
                current_token_count += paragraph_token_count
            else:
                # Otherwise, start a new chunk
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = [paragraphs[i]]
                current_token_count = paragraph_token_count
                current_idx = i
        
        # Add the final chunk
        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))
        
        return chunks
    except Exception as e:
        # If anything goes wrong, fall back to token chunking
        st.warning(f"Semantic chunking error: {str(e)}. Falling back to token-based chunking.")
        return chunk_by_tokens(text, max_tokens)

# Function to chunk text based on structural elements (headings)
def chunk_by_structure(text, max_tokens=500):
    """Chunk text based on headings and document structure."""
    # Split text by headings
    heading_pattern = r'(^|\n)(#+\s+.+)(\n|$)'
    parts = re.split(heading_pattern, text, flags=re.MULTILINE)
    
    # Combine the split parts back together
    sections = []
    for i in range(1, len(parts), 4):
        if i+2 < len(parts):
            # Combine the heading with the content that follows
            heading = parts[i+1]
            content = parts[i+2]
            sections.append(heading + content)
    
    # If no headings were found, fall back to paragraph chunking
    if not sections:
        return chunk_by_paragraphs(text, max_paragraphs=max(3, max_tokens // 100))
    
    # Process each section to ensure it doesn't exceed max_tokens
    chunks = []
    for section in sections:
        section_tokens = len(word_tokenize(section))
        if section_tokens <= max_tokens:
            chunks.append(section)
        else:
            # If the section is too large, split it further
            sub_chunks = chunk_by_tokens(section, max_tokens)
            chunks.extend(sub_chunks)
    
    return chunks

# Function to chunk text based on approximate token count (original implementation)
def chunk_by_tokens(text, max_tokens=500):
    try:
        sentences = sent_tokenize(text)
    except LookupError:
        # If punkt_tab is missing, fall back to a simpler split
        nltk.download('punkt', quiet=True)
        try:
            sentences = sent_tokenize(text)
        except:
            # Last resort: split by periods, question marks, and exclamation points
            sentences = re.split(r'(?<=[.!?])\s+', text)
    
    chunks = []
    current_chunk = []
    current_token_count = 0
    
    for sentence in sentences:
        # Approximate token count by words (not perfect but simple)
        try:
            sentence_token_count = len(word_tokenize(sentence))
        except:
            # Fall back to simple word counting if tokenization fails
            sentence_token_count = len(sentence.split())
        
        if current_token_count + sentence_token_count > max_tokens and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_token_count = sentence_token_count
        else:
            current_chunk.append(sentence)
            current_token_count += sentence_token_count
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

# Function to extract keywords using TF-IDF-like approach
def extract_keywords(text, max_keywords=10):
    try:
        # Try to load stopwords, download if not available
        try:
            stops = set(stopwords.words('english'))
        except LookupError:
            nltk.download('stopwords', quiet=True)
            stops = set(stopwords.words('english'))
        
        # Add punctuation to stopwords
        stops.update(set(punctuation))
        stops.update({'•', '■', '►', '●', '◆', '»', '"', '"', ''', ''', '–', '—'})
        
        # Tokenize and clean words
        words = word_tokenize(text.lower())
        words = [word for word in words if word not in stops and len(word) > 2]
        
        # Count word frequencies
        word_freq = Counter(words)
        
        # Get the most common words
        keywords = word_freq.most_common(max_keywords)
        
        return [word for word, count in keywords]
    
    except Exception as e:
        return [f"Error extracting keywords: {str(e)}"]

# Function to add metadata to chunks
def add_metadata(chunks, filename, chunking_method):
    result = []
    for i, chunk in enumerate(chunks):
        # Extract keywords for this chunk if enabled
        keywords_str = ""
        if 'extract_keywords' in st.session_state and st.session_state.extract_keywords:
            max_kw = st.session_state.max_keywords if 'max_keywords' in st.session_state else 10
            keywords = extract_keywords(chunk, max_kw)
            keywords_str = ", ".join(keywords)
            keywords_section = f"KEYWORDS: {keywords_str}\n"
        else:
            keywords_section = ""
        
        # Calculate approximate token count for reference
        token_count = len(word_tokenize(chunk))
        
        metadata = f"""---
CHUNK: {i+1}/{len(chunks)}
SOURCE: {filename}
CHUNKING_METHOD: {chunking_method}
TOKENS: ~{token_count}
{keywords_section}---

"""
        result.append(metadata + chunk)
    
    return result

# Function to create download link
def create_download_link(content, filename):
    b64 = base64.b64encode(content.encode()).decode()
    href = f'<a href="data:file/txt;base64,{b64}" download="{filename}">Download {filename}</a>'
    return href

# Function to create zip download link
def create_zip_download_link(file_contents, filenames, zip_filename):
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'a', zipfile.ZIP_DEFLATED) as zip_file:
        for content, filename in zip(file_contents, filenames):
            zip_file.writestr(filename, content)
    
    zip_buffer.seek(0)
    b64 = base64.b64encode(zip_buffer.getvalue()).decode()
    href = f'<a href="data:application/zip;base64,{b64}" download="{zip_filename}">Download {zip_filename}</a>'
    return href

# Sidebar navigation
with st.sidebar:
    st.markdown('<div class="sidebar-header">FileMerger App</div>', unsafe_allow_html=True)
    st.image("https://img.icons8.com/fluency/96/documents.png", width=80)
    
    # Navigation
    st.markdown("## Navigation")
    if st.button("Home", use_container_width=True):
        st.session_state.current_page = "home"
    if st.button("Upload Files", use_container_width=True):
        st.session_state.current_page = "upload"
    if st.button("Process Files", use_container_width=True):
        st.session_state.current_page = "process"
    if st.button("About", use_container_width=True):
        st.session_state.current_page = "about"
    
    st.markdown("---")
    st.markdown("### Options")
    
    # Only show options when in process page
    if st.session_state.current_page == "process":
        st.session_state.output_format = st.radio(
            "Output Format",
            ["Single Merged File", "Multiple Files (ZIP)"]
        )
        
        st.markdown("---")
        st.markdown("### Chunking Options")
        st.session_state.chunking_enabled = st.checkbox("Enable Chunking", value=False)
        
        if st.session_state.chunking_enabled:
            st.session_state.chunking_method = st.radio(
                "Chunking Method",
                ["Smart Paragraph", "Semantic", "Structure-based", "Token-based"]
            )
            
            if st.session_state.chunking_method == "Smart Paragraph":
                st.session_state.paragraphs_per_chunk = st.slider(
                    "Max Paragraphs per Chunk",
                    min_value=1,
                    max_value=20,
                    value=5
                )
                chunking_help = """
                Smart paragraph chunking keeps related paragraphs together and creates new chunks at logical boundaries such as headings.
                """
                
            elif st.session_state.chunking_method == "Semantic":
                st.session_state.tokens_per_chunk = st.slider(
                    "Max Tokens per Chunk",
                    min_value=100,
                    max_value=2000,
                    value=500,
                    step=100
                )
                st.session_state.similarity_threshold = st.slider(
                    "Similarity Threshold",
                    min_value=0.1,
                    max_value=0.5,
                    value=0.3,
                    step=0.05,
                    help="Higher values create more chunks with tighter semantic coherence"
                )
                chunking_help = """
                Semantic chunking groups paragraphs based on their content similarity, requiring scikit-learn package.
                """
                if not has_sklearn:
                    st.warning("Semantic chunking requires scikit-learn. It will fall back to token-based chunking.")
                
            elif st.session_state.chunking_method == "Structure-based":
                st.session_state.tokens_per_chunk = st.slider(
                    "Max Tokens per Chunk",
                    min_value=100,
                    max_value=2000,
                    value=500,
                    step=100
                )
                chunking_help = """
                Structure-based chunking creates chunks based on document headings and sections.
                """
                
            else:  # Token-based
                st.session_state.tokens_per_chunk = st.slider(
                    "Approximate Tokens per Chunk",
                    min_value=100,
                    max_value=2000,
                    value=500,
                    step=100
                )
                chunking_help = """
                Token-based chunking aims to keep chunks under the specified token count, splitting at sentence boundaries.
                """
            
            st.info(chunking_help)
        
        st.markdown("---")
        st.markdown("### Metadata Options")
        st.session_state.add_metadata = st.checkbox("Add Metadata", value=True)
        
        if st.session_state.add_metadata:
            st.session_state.extract_keywords = st.checkbox("Extract Keywords", value=True)
            
            if st.session_state.extract_keywords:
                st.session_state.max_keywords = st.slider(
                    "Maximum Keywords per Chunk",
                    min_value=3,
                    max_value=20,
                    value=10
                )

# Main area content based on current page
if st.session_state.current_page == "home":
    st.markdown('<div class="main-header">FileMerger App</div>', unsafe_allow_html=True)
    
    st.markdown("""
    Convert your documents to Markdown format for easy processing by AI tools.
    
    ### Features:
    - Convert PDF, Word, Excel, PowerPoint, HTML and text files to Markdown
    - Merge multiple files or download individually
    - Chunk documents for better AI processing
    - Add metadata for improved context
    
    ### How to use:
    1. Go to **Upload Files** and select your documents
    2. Navigate to **Process Files** to convert and customize
    3. Choose your output format and chunking options
    4. Download the processed files
    """)
    
    st.markdown('<div class="sub-header">Get Started</div>', unsafe_allow_html=True)
    if st.button("Go to Upload Files", type="primary"):
        st.session_state.current_page = "upload"

elif st.session_state.current_page == "upload":
    st.markdown('<div class="main-header">Upload Files</div>', unsafe_allow_html=True)
    
    # File uploader with extended file types
    uploaded_files = st.file_uploader(
        "Upload your documents",
        type=["pdf", "docx", "doc", "txt", "md", "csv", "xlsx", "xls", "pptx", "ppt", 
              "html", "htm", "json", "xml", "yaml", "yml", "text", "markdown"],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        st.session_state.uploaded_files = uploaded_files
        
        st.markdown('<div class="sub-header">Uploaded Files</div>', unsafe_allow_html=True)
        
        for file in uploaded_files:
            file_size = len(file.getvalue()) / 1024  # Size in KB
            
            st.markdown(f"""
            <div class="file-info">
                <strong>{file.name}</strong><br>
                Size: {file_size:.2f} KB<br>
                Type: {file.type}
            </div>
            """, unsafe_allow_html=True)
        
        if st.button("Proceed to Processing", type="primary"):
            st.session_state.current_page = "process"
    else:
        st.info("Please upload one or more files to continue.")

elif st.session_state.current_page == "process":
    st.markdown('<div class="main-header">Process Files</div>', unsafe_allow_html=True)
    
    if not st.session_state.uploaded_files:
        st.warning("No files uploaded. Please upload files first.")
        if st.button("Go to Upload"):
            st.session_state.current_page = "upload"
    else:
        # Show processing options
        st.markdown('<div class="sub-header">Processing Options</div>', unsafe_allow_html=True)
        st.write("Adjust the options in the sidebar, then click the button below to process your files.")
        
        if st.button("Process Files", type="primary"):
            with st.spinner("Processing files..."):
                processed_contents = []
                filenames = []
                
                # Process each file
                progress_bar = st.progress(0)
                for i, file in enumerate(st.session_state.uploaded_files):
                    # Update progress
                    progress_bar.progress((i+1)/len(st.session_state.uploaded_files))
                    
                    # Process file based on type
                    file_extension = os.path.splitext(file.name)[1]
                    content = process_file(file, file_extension)
                    
                    # Apply chunking if enabled
                    if st.session_state.chunking_enabled:
                        if st.session_state.chunking_method == "Smart Paragraph":
                            chunks = chunk_by_paragraphs(content, st.session_state.paragraphs_per_chunk)
                            chunking_method = f"Smart Paragraph ({st.session_state.paragraphs_per_chunk} max paragraphs per chunk)"
                        elif st.session_state.chunking_method == "Semantic":
                            similarity = st.session_state.similarity_threshold if 'similarity_threshold' in st.session_state else 0.3
                            chunks = chunk_by_semantic_similarity(content, st.session_state.tokens_per_chunk, similarity)
                            chunking_method = f"Semantic (~{st.session_state.tokens_per_chunk} tokens per chunk, similarity: {similarity})"
                        elif st.session_state.chunking_method == "Structure-based":
                            chunks = chunk_by_structure(content, st.session_state.tokens_per_chunk)
                            chunking_method = f"Structure-based (~{st.session_state.tokens_per_chunk} tokens per chunk)"
                        else:
                            chunks = chunk_by_tokens(content, st.session_state.tokens_per_chunk)
                            chunking_method = f"Token-based (~{st.session_state.tokens_per_chunk} tokens per chunk)"
                        
                        # Add metadata if enabled
                        if st.session_state.add_metadata:
                            chunks = add_metadata(chunks, file.name, chunking_method)
                        
                        # For individual files
                        processed_content = "\n\n".join(chunks)
                    else:
                        # If chunking is disabled, just add a simple header
                        if st.session_state.add_metadata:
                            processed_content = f"""---
SOURCE: {file.name}
---

{content}"""
                        else:
                            processed_content = content
                    
                    processed_contents.append(processed_content)
                    filenames.append(f"{os.path.splitext(file.name)[0]}.md")
                
                # Store processed files in session state
                st.session_state.processed_files = list(zip(processed_contents, filenames))
                
                # Create merged content for single file option
                if st.session_state.add_metadata:
                    merged_content = "\n\n---\n\n".join(processed_contents)
                else:
                    merged_content = "\n\n".join(processed_contents)
                
                st.session_state.merged_content = merged_content
                st.session_state.conversion_done = True
            
            # Clear progress bar after completion
            progress_bar.empty()
            st.success("Processing complete!")
        
        # Show download options if processing is done
        if st.session_state.conversion_done:
            st.markdown('<div class="sub-header">Download Results</div>', unsafe_allow_html=True)
            
            # Preview tab
            with st.expander("Preview Converted Content", expanded=True):
                if st.session_state.output_format == "Single Merged File":
                    st.text_area("Merged Content Preview", st.session_state.merged_content[:1000] + "...", height=300)
                else:
                    if st.session_state.processed_files:
                        sample_file = st.session_state.processed_files[0]
                        st.text_area(f"Preview of {sample_file[1]}", sample_file[0][:1000] + "...", height=300)
                        
                        if len(st.session_state.processed_files) > 1:
                            st.info(f"{len(st.session_state.processed_files)-1} more files are available for download.")
            
            # Download options
            col1, col2 = st.columns(2)
            with col1:
                if st.session_state.output_format == "Single Merged File":
                    st.markdown(create_download_link(st.session_state.merged_content, "merged_document.md"), unsafe_allow_html=True)
                else:
                    contents, names = zip(*st.session_state.processed_files)
                    st.markdown(create_zip_download_link(contents, names, "converted_files.zip"), unsafe_allow_html=True)
            
            with col2:
                if st.button("Process More Files"):
                    st.session_state.conversion_done = False
                    st.session_state.uploaded_files = []
                    st.session_state.current_page = "upload"

elif st.session_state.current_page == "about":
    st.markdown('<div class="main-header">About FileMerger App</div>', unsafe_allow_html=True)
    
    st.markdown("""
    **FileMerger App** is a tool designed to help you prepare documents for AI processing.
    
    ### Purpose
    This app converts various document formats to Markdown, which is easily consumable by AI tools and large language models. You can merge multiple documents or process them individually, with options for chunking to make the content more digestible.
    
    ### Supported Formats
    - PDF (.pdf)
    - Word (.docx, .doc)
    - Text (.txt, .text)
    - Markdown (.md, .markdown)
    - Excel/CSV (.xlsx, .xls, .csv)
    - PowerPoint (.pptx, .ppt)
    - HTML (.html, .htm)
    - Data formats (.json, .xml, .yaml, .yml)
    
    ### Privacy
    All processing happens in your browser session. Your files are not stored on any server and are only held in memory during processing.
    
    ### Feedback
    If you encounter any issues or have suggestions for improvements, please reach out.
    """)
    
    if st.button("Return to Home"):
        st.session_state.current_page = "home" 